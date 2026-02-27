# multi_format_to_json.py
import os
import json
from pathlib import Path
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import pandas as pd
from PIL import Image
import pytesseract
import io
import logging

# 可选：设置 Tesseract 路径（Windows 若未加 PATH 需手动指定）
# pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

def ocr_image(image: Image.Image, lang: str = 'chi_sim+eng') -> str:
    """对 PIL 图像进行 OCR，返回识别文本"""
    try:
        # 可选：预处理提升 OCR 准确率（灰度 + 二值化）
        gray = image.convert('L')
        # 简单二值化
        thresh = 128
        fn = lambda x: 255 if x > thresh else 0
        binary = gray.point(fn, mode='1')
        
        text = pytesseract.image_to_string(binary, lang=lang)
        return text.strip()
    except Exception as e:
        logger.warning(f"OCR 失败: {e}")
        return ""

# ========================
# 增强版提取函数（带图片 OCR）
# ========================

def extract_pdf(file_path: str) -> str:
    reader = PdfReader(file_path)
    text = ""
    for page_num, page in enumerate(reader.pages):
        # 1. 提取文本
        extracted = page.extract_text()
        if extracted:
            text += f"Page {page_num + 1}:\n{extracted}\n\n"
        
        # 2. 提取图片并 OCR（PyPDF2 不支持直接提图，跳过）
        # 注：PyPDF2 无法提取 PDF 中的图片！需用 pdf2image + poppler
        # 暂不实现（复杂），仅提示
        # TODO: 如需 PDF 图片 OCR，请改用 pdf2image + Tesseract
    return text.strip()

def extract_docx(file_path: str) -> str:
    doc = Document(file_path)
    text = ""
    
    # 1. 段落文本
    for para in doc.paragraphs:
        if para.text.strip():
            text += para.text + "\n"
    
    # 2. 表格
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                text += row_text + "\n"
    
    # 3. 图片（docx 中图片在 drawing 或 blip 中）
    image_count = 0
    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            try:
                image_data = rel.target_part.blob
                pil_img = Image.open(io.BytesIO(image_data))
                ocr_result = ocr_image(pil_img)
                if ocr_result:
                    text += f"[图片OCR #{image_count + 1}]:\n{ocr_result}\n\n"
                image_count += 1
            except Exception as e:
                logger.warning(f"DOCX 图片 OCR 失败: {e}")
    
    return text.strip()

def extract_xlsx(file_path: str) -> str:
    # Excel 一般不含图片（或极难提取），跳过 OCR
    excel_file = pd.ExcelFile(file_path)
    all_text = ""
    for sheet_name in excel_file.sheet_names:
        df = pd.read_excel(file_path, sheet_name=sheet_name, dtype=str)
        df = df.fillna("")
        sheet_text = f"Sheet: {sheet_name}\n"
        for _, row in df.iterrows():
            row_text = " | ".join(str(cell) for cell in row if str(cell).strip())
            if row_text.strip():
                sheet_text += row_text + "\n"
        all_text += sheet_text + "\n"
    return all_text.strip()

def extract_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    text = ""
    for slide_num, slide in enumerate(prs.slides):
        slide_text = f"Slide {slide_num + 1}:\n"
        
        # 1. 文本
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text += shape.text.strip() + "\n"
        
        # 2. 图片 OCR
        image_count = 0
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                try:
                    image = shape.image
                    image_bytes = image.blob
                    pil_img = Image.open(io.BytesIO(image_bytes))
                    ocr_result = ocr_image(pil_img)
                    if ocr_result:
                        slide_text += f"[图片OCR #{image_count + 1}]:\n{ocr_result}\n"
                    image_count += 1
                except Exception as e:
                    logger.warning(f"PPTX 图片 OCR 失败 (Slide {slide_num + 1}): {e}")
        
        if "Slide" in slide_text:
            text += slide_text + "\n"
    
    return text.strip()

# ========================
# 其余代码保持不变
# ========================

def extract_file(file_path: str) -> dict:
    file_path = str(file_path)
    ext = file_path.lower().split('.')[-1]
    
    try:
        if ext == 'pdf':
            content = extract_pdf(file_path)
            file_type = 'pdf'
        elif ext == 'docx':
            content = extract_docx(file_path)
            file_type = 'docx'
        elif ext == 'xlsx':
            content = extract_xlsx(file_path)
            file_type = 'xlsx'
        elif ext == 'pptx':
            content = extract_pptx(file_path)
            file_type = 'pptx'
        else:
            return None
        
        return {
            "filename": os.path.basename(file_path),
            "file_type": file_type,
            "content": content,
            "char_count": len(content),
            "source_path": file_path
        }
    except Exception as e:
        print(f"❌ 解析失败: {file_path} | 错误: {e}")
        return None

def batch_convert_to_json(input_dir: str, output_file: str):
    input_path = Path(input_dir)
    if not input_path.exists():
        raise FileNotFoundError(f"输入目录不存在: {input_dir}")
    
    results = []
    supported_exts = ('.pdf', '.docx', '.xlsx', '.pptx')
    
    for file_path in input_path.rglob('*'):
        if file_path.is_file() and file_path.suffix.lower() in supported_exts:
            print(f"📄 正在处理: {file_path.name}")
            doc = extract_file(str(file_path))
            if doc:
                results.append(doc)
    
    output_path = Path(output_file)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(results, f, ensure_ascii=False, indent=2)
    
    print(f"\n✅ 成功处理 {len(results)} 个文件")
    print(f"📁 输出文件: {output_path.absolute()}")

if __name__ == "__main__":
    INPUT_DIR = "input_docs"
    OUTPUT_FILE = "output/combined_docs.json"
    batch_convert_to_json(INPUT_DIR, OUTPUT_FILE)